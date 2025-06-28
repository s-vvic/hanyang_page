import os
import math
from collections import Counter
import numpy as np
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import openpyxl
#pip install python-docx python-pptx pymupdf openpyxl numpy <-- 이거 해야됨

# ------------------- 파일 읽기 함수들 ----------------------

def read_txt(path):             #txt 파일 읽는 함수
    with open(path, "r", encoding="utf-8") as f:  # "r" 모드는 읽기 모드 "encoding" 은 한글이 깨지지 않게 열기
        return f.read()                           # "with"는 파일을 자동으로 열고 닫기    

def read_docx(path):            #docx(워드) 파일 읽는 함수
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def read_pdf(path):             #pdf 파일 읽는 함수
    try:
        doc = fitz.open(path)
        return "\n".join(page.get_text() for page in doc)
    except Exception as e:
        print(f"[PDF 오류] {path} - {e}")
        return ""

def read_xlsx(path):            #엑셀 파일 읽는 함수
    wb = openpyxl.load_workbook(path, data_only=True)
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value:
                    text += str(cell.value) + " "
    return text

def read_hwp_txt(path):         #hwp(한글) 파일 읽는 함수
    return read_txt(path)       #한글 파일은 읽기 빡세서 미리 텍스트 파일로 변환해서 읽기

def read_pptx(path):            #ppt 파일 읽는 함수
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def read_any_file(path):        #종합
    if path.endswith(".txt") and not path.endswith(".hwp.txt"):
        return read_txt(path)
    elif path.endswith(".docx"):
        return read_docx(path)
    elif path.endswith(".pdf"):
        return read_pdf(path)
    elif path.endswith(".xlsx"):
        return read_xlsx(path)
    elif path.endswith(".hwp.txt"):
        return read_hwp_txt(path)
    elif path.endswith(".pptx") or path.endswith(".ppt"):
        return read_pptx(path)
    else:
        return ""

def load_documents_from_folder(folder_path):    #폴더 안에 있는 문서를 읽어서 내용이랑 파일 이름을 한번에 가져오는 함수
    docs = []                                   #문서 내용들을 저장하는 리스트
    filenames = []                              #각 문서의 파일 이름을 저장할 리스트
    for filename in os.listdir(folder_path):    #폴더에 있는 모든 파일들의 이름을 리스트로 반환
        full_path = os.path.join(folder_path, filename)     # 전체 경로 생성
        if os.path.isfile(full_path):           #진짜 파일인지 확인하는 조건문
            content = read_any_file(full_path)  #read_any_file() 함수로 파일을 자동 형식 감지해서 읽음
            if content.strip():
                docs.append(content)            #문서 내용을 docs 리스트에 추가
                filenames.append(filename)      #해당 문서의 파일명을 filenames 리스트에 추가
    return docs, filenames

#-----------------------------------------------------

def tokenize_document(doc_text):
    return doc_text.lower().split()

def compute_idf(corpus, k1=1.5, b=0.75):
    num_documents = len(corpus)
    doc_freq = Counter()
    for doc in corpus:
        unique_terms_in_doc = set(doc)
        for term in unique_terms_in_doc:
            doc_freq[term] += 1

    idf_scores = {}
    for term, df in doc_freq.items():
        idf_scores[term] = math.log((num_documents - df + 0.5) / (df + 0.5) + 1)
    return idf_scores, doc_freq

def bm25_score(query, document, idf_scores, avgdl, k1=1.5, b=0.75):
    score = 0.0
    doc_len = len(document)
    term_counts = Counter(document)

    for query_term in query:
        if query_term not in idf_scores:
            continue

        idf = idf_scores[query_term]
        tf = term_counts[query_term]

        numerator = tf * (k1 + 1)
        denominator = tf + k1 * (1 - b + b * (doc_len / avgdl))

        score += idf * (numerator / denominator)
    return score

# ---------------------- 실행 시작 ----------------------

folder_path = "./documents"
raw_corpus, filenames = load_documents_from_folder(folder_path)

# 1. 문서 토큰화
tokenized_corpus = [tokenize_document(doc) for doc in raw_corpus]

# 2. 문서 길이 및 평균 길이 계산
document_lengths = [len(doc) for doc in tokenized_corpus]
avgdl = sum(document_lengths) / len(document_lengths)

# 3. IDF 값 계산
idf_scores, _ = compute_idf(tokenized_corpus)

# 4. 쿼리 준비
query_text = input("검색어를 입력하세요: ")
tokenized_query = tokenize_document(query_text)

# 5. 각 문서에 대한 BM25 점수 계산
scores = []
for i, doc in enumerate(tokenized_corpus):
    score = bm25_score(tokenized_query, doc, idf_scores, avgdl)
    scores.append(score)
    print(f"문서 {i} ('{filenames[i]}'): {score:.4f}")

# 6. 결과 랭킹
ranked_indices = np.argsort(scores)[::-1]
print("\n랭킹된 문서 인덱스 (높은 점수부터):", ranked_indices)

print("\n랭킹된 문서 내용:")
for idx in ranked_indices:
    print(f"문서 {idx}: {filenames[idx]} (점수: {scores[idx]:.4f})")
