#------------------------------------유저 입력 받는 코드------------------------------------------

class Get_User_Query:
    def __init__(self):
        self.query = None # 초기에는 쿼리가 없으므로 None으로 설정
        self.query_history = [] # 쿼리 이력을 저장하는 리스트 / 필요 없으면 지워도 됨

    def get_user_query(self):
        while True: # 유효한 쿼리를 받을 때까지 반복
            query_input = input("검색할 내용을 입력하세요 (최소 2자 이상): ").strip() # 공백 제거
            if len(query_input) >= 2: # 최소 2자 이상인지 확인
                self.query = query_input.lower()

                if self.query in self.query_history:
                    self.query_history.remove(self.query)

                if self.query:
                    self.query_history.append(self.query)

                print(f"입력된 검색 쿼리: '{self.query}'")
                print(f"최근 검색 이력 (최신순 정렬): {self.query_history[::-1]}")
                return self.query
            else:
                print("검색 쿼리는 최소 2자 이상이어야 합니다. 다시 입력해주세요.")

#----------------------------------------------------------------------------------------------

#--------------------------------폴더 내 모든 문서 읽는 코드---------------------------------------
    
import os
from docx import Document
from pptx import Presentation
import fitz
import openpyxl

class Search_docs:
    def __init__(self):
        self.folder_path = "./documents"
    
    def read_txt(self, path):             #txt 파일 읽는 함수
        with open(path, "r", encoding="utf-8") as f:  # "r" 모드는 읽기 모드 "encoding" 은 한글이 깨지지 않게 열기
            return f.read()                           # "with"는 파일을 자동으로 열고 닫기    

    def read_docx(self, path):            #docx(워드) 파일 읽는 함수
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)

    def read_pdf(self, path):             #pdf 파일 읽는 함수
        try:
            doc = fitz.open(path)
            return "\n".join(page.get_text() for page in doc)
        except Exception as e:
            print(f"[PDF 오류] {path} - {e}")
            return ""

    def read_xlsx(self, path):            #엑셀 파일 읽는 함수
        wb = openpyxl.load_workbook(path, data_only=True)
        text = ""
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        text += str(cell.value) + " "
        return text

    def read_hwp_txt(self, path):         #hwp(한글) 파일 읽는 함수
        return self.read_txt(path)       #한글 파일은 읽기 빡세서 미리 텍스트 파일로 변환해서 읽기

    def read_pptx(self, path):            #ppt 파일 읽는 함수
        prs = Presentation(path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    def read_any_file(self, path):        #종합
        if path.endswith(".txt") and not path.endswith(".hwp.txt"):
            return self.read_txt(path)
        elif path.endswith(".docx"):
            return self.read_docx(path)
        elif path.endswith(".pdf"):
            return self.read_pdf(path)
        elif path.endswith(".xlsx"):
            return self.read_xlsx(path)
        elif path.endswith(".hwp.txt"):
            return self.read_hwp_txt(path)
        elif path.endswith(".pptx") or path.endswith(".ppt"):
            return self.read_pptx(path)
        else:
            return ""

    def load_documents_from_folder(self):    #폴더 안에 있는 문서를 읽어서 내용이랑 파일 이름을 한번에 가져오는 함수
        docs = []                                   #문서 내용들을 저장하는 리스트
        filenames = []                              #각 문서의 파일 이름을 저장할 리스트
        for filename in os.listdir(self.folder_path):    #폴더에 있는 모든 파일들의 이름을 리스트로 반환
            full_path = os.path.join(self.folder_path, filename)     # 전체 경로 생성
            if os.path.isfile(full_path):           #진짜 파일인지 확인하는 조건문
                content = self.read_any_file(full_path)  #read_any_file() 함수로 파일을 자동 형식 감지해서 읽음
                if content.strip():
                    docs.append(content)            #문서 내용을 docs 리스트에 추가
                    filenames.append(filename)      #해당 문서의 파일명을 filenames 리스트에 추가
        return docs, filenames
    
#----------------------------------------------------------------------------------------------

#-------------------------------문서의 토큰화, 불용어 처리 코드------------------------------------
    
from konlpy.tag import Okt
from nltk.tokenize import word_tokenize
from nltk import pos_tag
from nltk.corpus import stopwords

class Preprocessing:
    # nltk 리소스 다운로드
    #nltk.download('punkt')
    #nltk.download('averaged_perceptron_tagger')
    #nltk.download('averaged_perceptron_tagger_eng')

    # 한글 문서인지 영어 문서인지 판단
    def isEnglishOrKorean(self, input_s):
        k_count = 0
        e_count = 0
        for c in input_s:
            if ord('가') <= ord(c) <= ord('힣'):
                k_count+=1
            elif ord('a') <= ord(c.lower()) <= ord('z'):
                e_count+=1
        return 1 if k_count>e_count else 0

    def Tokenize(self, docs):
        for text in docs:
            # 한글 문서라면
            if self.isEnglishOrKorean(text) == 1:

                # Okt 형태소 분석기 인스턴스 생성
                okt = Okt()

                # 명사 추출
                nouns = okt.nouns(text)
                return nouns

            # 영어 문서라면
            elif self.isEnglishOrKorean(text) == 0:
                # 토큰화 (형태소 분석)
                tokens = word_tokenize(text)
                #print("Tokens:", tokens)

                # 품사 태깅
                pos_tags = pos_tag(tokens)
                #print("POS Tags:", pos_tags)

                # 명사 추출 (품사 태그가 NN, NNS, NNP, NNPS인 단어 추출)
                nouns = [word for word, pos in pos_tags if pos in ['NN', 'NNS', 'NNP', 'NNPS']]

                return nouns

#----------------------------------------------------------------------------------------------

# ------------------ query와 문서 내용의 유사도 계산 및 관련도 높은 문서 랭킹 출력 --------------------

import math
from collections import Counter

class BM_25:
    def tokenize_document(doc_text):
        return doc_text.lower().split()

    def compute_idf(corpus, k1=1.5, b=0.75):
        num_documents = len(corpus)
        doc_freq = Counter()
        for doc in corpus:
            unique_terms_in_doc = set(doc)
            for term in unique_terms_in_doc:
                doc_freq[term] += 1

        idf_scores = {}
        for term, df in doc_freq.items():
            idf_scores[term] = math.log((num_documents - df + 0.5) / (df + 0.5) + 1)
        return idf_scores, doc_freq

    def bm25_score(query, document, idf_scores, avgdl, k1=1.5, b=0.75):
        score = 0.0
        doc_len = len(document)
        term_counts = Counter(document)

        for query_term in query:
            if query_term not in idf_scores:
                continue

            idf = idf_scores[query_term]
            tf = term_counts[query_term]

            numerator = tf * (k1 + 1)
            denominator = tf + k1 * (1 - b + b * (doc_len / avgdl))

            score += idf * (numerator / denominator)
        return score